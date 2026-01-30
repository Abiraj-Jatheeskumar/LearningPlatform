import io
from typing import List
from fastapi import UploadFile
import PyPDF2


async def extract_text_from_pdf(file: UploadFile) -> List[str]:
    """
    Extract text from PDF file, one item per page
    """
    try:
        content = await file.read()
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        texts = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text()
            if text and len(text.strip()) > 20:
                texts.append(text.strip())
        
        return texts
    except Exception as e:
        raise Exception(f"Failed to extract PDF text: {str(e)}")


async def extract_text_from_pptx(file: UploadFile) -> List[str]:
    """
    Extract text from PowerPoint file, one item per slide
    """
    try:
        from pptx import Presentation
        
        content = await file.read()
        ppt_file = io.BytesIO(content)
        presentation = Presentation(ppt_file)
        
        texts = []
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            full_text = "\n".join(slide_text).strip()
            if full_text and len(full_text) > 20:
                texts.append(full_text)
        
        return texts
    except Exception as e:
        raise Exception(f"Failed to extract PPTX text: {str(e)}")


async def extract_text_from_file(file: UploadFile) -> List[str]:
    """
    Extract text from uploaded file based on file type
    """
    filename = file.filename.lower()
    
    if filename.endswith('.pdf'):
        return await extract_text_from_pdf(file)
    elif filename.endswith('.pptx'):
        return await extract_text_from_pptx(file)
    else:
        raise ValueError(f"Unsupported file type: {filename}. Only PDF and PPTX are supported.")
